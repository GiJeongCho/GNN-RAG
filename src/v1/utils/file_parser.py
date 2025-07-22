# src/utils/file_parser.py
"""Utility for converting many document types to plain text strings.

Supported extensions are declared in SUPPORTED_EXTENSIONS. Each parser tries
its best using optional third-party libraries; when a required library is
missing it raises RuntimeError so the caller can handle the failure.
"""
from __future__ import annotations

import io
import csv
from pathlib import Path
from typing import Optional, Set

# ----------------------------
# Public API
# ----------------------------

SUPPORTED_EXTENSIONS: Set[str] = {
    ".mbox", ".pdf", ".html", ".htm", ".txt", ".md", ".org", ".adoc", ".rst",
    ".docx", ".pptx", ".odt", ".odp", ".xlsx", ".csv", ".hwp",
}

__all__ = ["SUPPORTED_EXTENSIONS", "parse_file"]

# ----------------------------
# Helpers
# ----------------------------

def _read_plain_text(path: Path, encoding: str = "utf-8") -> str:
    """Read a file as plain text using best-guess encoding."""
    try:
        return path.read_text(encoding=encoding)
    except UnicodeDecodeError:
        try:
            import chardet  # type: ignore
        except ImportError as exc:
            raise RuntimeError("chardet not installed for encoding detection") from exc
        raw = path.read_bytes()
        det = chardet.detect(raw)
        enc = det["encoding"] or "utf-8"
        return raw.decode(enc, errors="ignore")

# ----------------------------
# Individual parsers
# ----------------------------

def _parse_pdf(path: Path) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    text_pages = [p.get_text("text", sort=True).strip() for p in doc]
    text = "\n\n".join(text_pages)

    # OCR image blocks if libraries are available
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore

        for page_index in range(len(doc)):
            for img in doc.get_page_images(page_index):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image.get("image")
                if img_bytes:
                    img_text = pytesseract.image_to_string(
                        Image.open(io.BytesIO(img_bytes)), lang="kor+eng"
                    ).strip()
                    if img_text:
                        text += "\n\n" + img_text
    except Exception:
        # OCR is best-effort; ignore all failures.
        pass

    return text


def _parse_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-docx not installed") from exc
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed") from exc
    prs = Presentation(path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _parse_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError as exc:
        raise RuntimeError("beautifulsoup4 not installed") from exc
    soup = BeautifulSoup(_read_plain_text(path), "html.parser")
    return soup.get_text(separator="\n")


def _parse_csv(path: Path) -> str:
    lines: list[str] = []
    with path.open(newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append(", ".join(row))
    return "\n".join(lines)


def _parse_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except ImportError as exc:
        raise RuntimeError("openpyxl not installed") from exc
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            lines.append(", ".join([str(cell) if cell is not None else "" for cell in row]))
    return "\n".join(lines)


def _parse_mbox(path: Path) -> str:
    import mailbox

    mbox = mailbox.mbox(path)
    lines: list[str] = []
    for msg in mbox:
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        lines.append(payload.decode("utf-8", errors="ignore"))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                lines.append(payload.decode("utf-8", errors="ignore"))
    return "\n".join(lines)


def _parse_hwp(path: Path) -> str:
    try:
        import pyhwp  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pyhwp not installed") from exc

    doc = pyhwp.HWPDocument(path)
    # pyhwp extracts text per section; join with newline
    return "\n".join(doc.body_text())

# ----------------------------
# Dispatcher
# ----------------------------

def parse_file(path: Path) -> Optional[str]:
    """Return plain text extracted from *path* or None when unsupported.

    Raises RuntimeError when the format is supported but required libraries are
    missing or the parsing fails for another reason. Caller should handle and
    log the exception.
    """
    ext = path.suffix.lower()

    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        if ext in {".txt", ".md", ".org", ".adoc", ".rst"}:
            return _read_plain_text(path)
        if ext in {".html", ".htm"}:
            return _parse_html(path)
        if ext == ".docx":
            return _parse_docx(path)
        if ext == ".pptx":
            return _parse_pptx(path)
        if ext in {".csv"}:
            return _parse_csv(path)
        if ext in {".xlsx"}:
            return _parse_xlsx(path)
        if ext == ".mbox":
            return _parse_mbox(path)
        if ext == ".hwp":
            return _parse_hwp(path)
        if ext in {".odt", ".odp"}:
            # Best-effort: treat as plain text (may contain XML)
            return _read_plain_text(path)
    except RuntimeError:
        raise
    except Exception as exc:
        raise RuntimeError(f"{ext} parsing failed: {exc}") from exc

    # Unsupported extension
    return None 

